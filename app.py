# app.py — Media Impact Analyzer (Quick ITS + Paper Study Set)
# Python 3.13 compatible. Two tabs:
#  1) New Title (Quick ITS): Google Trends & Wikipedia weekly signals -> segmented ITS + visuals
#  2) Study set (Paper): auto-fetch CSV outputs from the official repo results/ (or upload) -> tables & forest plots

import io, time, zipfile, datetime as dt
from typing import List, Optional, Tuple, Dict
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import statsmodels.api as sm
import streamlit as st
import requests
from urllib.parse import quote

# --------------------------- Page setup & light CSS ---------------------------
st.set_page_config(page_title="🎬 Media Impact Analyzer — Google Trends & Wikipedia", layout="wide")
st.markdown(
    """
    <style>
      .center-card {max-width: 920px; margin: 0 auto;}
      .stButton>button {width: 100%;}
      .muted {color: #9aa0a6;}
    </style>
    """, unsafe_allow_html=True
)

st.title("🎬 Media Impact Analyzer — Google Trends & Wikipedia")
st.caption(
    "Fetches weekly public–interest signals from Google Trends and the Wikimedia Pageviews API for your project "
    "title/outcome, then fits an Interrupted Time Series (level & slope change) at your release date. "
    "A second tab loads the paper’s prepared results (tables, forest plots, benchmarking)."
)

# =====================================================
# ---------- Data fetchers (robust, 3.13-safe) --------
# =====================================================
def fetch_trends_weekly(queries: List[str], geo: str, timeframe: str = "today 5-y",
                        hl: str = "en-US", tz: int = 0, tries: int = 6) -> pd.Series:
    """
    Google Trends weekly series (W-MON), averaged across up to 5 queries.
    Retries politely for 429s and widens timeframe if needed.
    """
    from pytrends.request import TrendReq
    from pytrends import exceptions as pex

    queries = [q.strip() for q in queries if q and q.strip()]
    if not queries:
        raise ValueError("Please provide at least one query.")
    queries = queries[:5]

    py = TrendReq(hl=hl, tz=tz, retries=0, backoff_factor=0, timeout=(10, 30),
                  requests_args={"headers": {"User-Agent": "media-impacts-app/0.2"}})

    wait = 5
    last_err = None
    for attempt in range(tries):
        try:
            py.build_payload(queries, timeframe=timeframe, geo=geo)
            df = py.interest_over_time()
            if df is None or df.empty:
                raise RuntimeError("Google Trends returned no data")
            if "isPartial" in df.columns:
                df = df.drop(columns=["isPartial"])
            s = df[queries].mean(axis=1).resample("W-MON").mean()
            if s is None or s.empty:
                raise RuntimeError("Google Trends series empty after resample")
            return s.rename("value")
        except Exception as e:
            last_err = e
            msg = str(e)
            if isinstance(e, pex.TooManyRequestsError) or "429" in msg:
                st.info(f"Google Trends rate limited. Retrying in {wait}s… (try {attempt+1}/{tries})")
                time.sleep(wait)
                wait = min(wait * 2, 60)
                continue
            if "code 400" in msg and timeframe != "today 5-y":
                timeframe = "today 5-y"
                continue
            if any(k in msg for k in ["502", "503", "504", "timeout"]):
                time.sleep(3)
                continue
            break
    raise RuntimeError(f"Google Trends failed after retries: {last_err}")

def wiki_weekly_pageviews(title_or_query: str, start: str, end: str, lang: str = "en") -> pd.Series:
    """
    Wikipedia weekly pageviews (sum of daily), W-MON index, with search fallback & friendly errors.
    """
    session = requests.Session()
    session.headers.update({"User-Agent": "media-impacts-app/0.2 (contact: you@example.com)"})

    def _summary_exists(title: str) -> bool:
        r = session.get(f"https://{lang}.wikipedia.org/api/rest_v1/page/summary/{quote(title, safe='')}", timeout=15)
        if r.status_code == 404:
            return False
        r.raise_for_status()
        return True

    title = title_or_query.strip()
    try:
        if not _summary_exists(title):
            sr = session.get(
                f"https://{lang}.wikipedia.org/w/api.php",
                params={"action": "query", "list": "search", "srsearch": title_or_query, "format": "json"},
                timeout=20,
            )
            sr.raise_for_status()
            hits = sr.json().get("query", {}).get("search", [])
            if not hits:
                raise RuntimeError(f"No Wikipedia page found for “{title_or_query}”.")
            title = hits[0]["title"]

        r = session.get(
            f"https://wikimedia.org/api/rest_v1/metrics/pageviews/per-article/"
            f"{lang}.wikipedia/all-access/user/{quote(title, safe='')}/daily/{start}/{end}",
            timeout=30,
        )
        r.raise_for_status()
    except requests.HTTPError as e:
        code = getattr(e.response, "status_code", "")
        raise RuntimeError(f"Wikipedia API error {code}. Please rerun or try again in a minute.") from e

    items = r.json().get("items", [])
    if not items:
        raise RuntimeError(f"No pageview data for “{title}”.")
    idx = [pd.to_datetime(it["timestamp"][:8]) for it in items]
    vals = [it["views"] for it in items]
    s = pd.Series(vals, index=pd.DatetimeIndex(idx), name=title).resample("W-MON").sum()
    if s.empty:
        raise RuntimeError(f"Wikipedia returned no data after resample for “{title}”.")
    return s

# =====================================================
# --------------------- Modeling ----------------------
# =====================================================
def quick_its(outcome: pd.Series, film: pd.Series, intervention_date: str):
    """Segmented regression: level + slope change w/ HAC errors. Returns model, df, figs, metrics."""
    df = pd.concat(
        {"outcome": outcome.rename("outcome"), "film_interest": film.rename("film_interest")},
        axis=1
    ).dropna()
    if len(df) < 12:
        raise RuntimeError("Not enough weekly data after alignment. Try a wider timeframe or different source.")

    df["time"] = np.arange(len(df))
    t0 = pd.to_datetime(intervention_date)
    t0_monday = (t0 - pd.offsets.Week(weekday=0))
    df["post"] = (df.index >= t0_monday).astype(int)
    df["time_post"] = df["time"] * df["post"]
    df["film_lag1"] = df["film_interest"].shift(1).fillna(method="bfill")

    X = sm.add_constant(df[["time", "post", "time_post", "film_lag1"]])
    y = df["outcome"]
    model = sm.OLS(y, X).fit(cov_type="HAC", cov_kwds={"maxlags": 4})

    # Figures
    fig1 = plt.figure(figsize=(10, 4))
    plt.plot(df.index, df["outcome"], label="Outcome interest")
    plt.axvline(t0_monday, linestyle="--", label=f"Intervention: {t0.date()}")
    plt.plot(df.index, model.predict(X), label="Fitted (segmented regression)")
    plt.legend(); plt.title("Outcome vs intervention — quick ITS"); plt.tight_layout()

    fig2 = plt.figure(figsize=(10, 2.8))
    plt.plot(df.index, df["film_interest"], label="Film interest")
    plt.axvline(t0_monday, linestyle="--"); plt.title("Film interest"); plt.tight_layout()

    pre_mean = df.loc[df.index < t0_monday, "outcome"].mean()
    ci = model.conf_int()
    level = float(model.params.get("post", np.nan))
    slope = float(model.params.get("time_post", np.nan))
    level_lo, level_hi = [float(x) for x in ci.loc["post"]] if "post" in ci.index else (np.nan, np.nan)
    slope_lo, slope_hi = [float(x) for x in ci.loc["time_post"]] if "time_post" in ci.index else (np.nan, np.nan)

    metrics = {
        "level_change": level,
        "level_change_pct_of_pre": float(100 * level / pre_mean) if pre_mean else np.nan,
        "level_ci_lo": level_lo, "level_ci_hi": level_hi,
        "slope_change_per_week": slope,
        "slope_ci_lo": slope_lo, "slope_ci_hi": slope_hi,
        "t0_monday": t0_monday, "pre_mean": float(pre_mean) if pre_mean else np.nan
    }
    return model, df, (fig1, fig2), metrics

def trends_interest_by_region(query: str, geo: str = "", timeframe: str = "today 5-y") -> pd.DataFrame:
    from pytrends.request import TrendReq
    py = TrendReq(hl="en-US", tz=0, retries=0, backoff_factor=0,
                  requests_args={"headers": {"User-Agent": "media-impacts-app/0.2"}})
    py.build_payload([query], timeframe=timeframe, geo=geo)
    df = py.interest_by_region(resolution="COUNTRY" if geo == "" else "REGION",
                               inc_low_vol=True, inc_geo_code=True).reset_index()
    name_col = "geoName" if "geoName" in df.columns else df.columns[0]
    out = df.rename(columns={name_col: "region"})
    return out.sort_values(query, ascending=False).reset_index(drop=True)

# =====================================================
# ------------------ Cache wrappers -------------------
# =====================================================
@st.cache_data(ttl=3600, show_spinner=False)
def cached_trends(queries, geo, timeframe):
    s = fetch_trends_weekly(queries, geo=geo, timeframe=timeframe)
    if s is None or not isinstance(s, pd.Series) or s.empty:
        raise RuntimeError("No Google Trends data returned. Try simpler terms, different GEO, or Wikipedia.")
    return s

@st.cache_data(ttl=3600, show_spinner=False)
def cached_wiki(title, start, end, lang="en"):
    return wiki_weekly_pageviews(title, start, end, lang=lang)

@st.cache_data(ttl=1800, show_spinner=False)
def cached_regions(query, geo, timeframe):
    return trends_interest_by_region(query, geo=geo, timeframe=timeframe)

# =====================================================
# ----------- Combined Trends call (pair) -------------
# =====================================================
def trends_pair(film_qs, outcome_qs, geo, timeframe):
    """Try film+outcome in one Trends request (≤5 total). Auto-widen timeframe/worldwide; fallback to separate calls."""
    from pytrends.request import TrendReq
    all_q = [q.strip() for q in (film_qs + outcome_qs) if q.strip()]
    py = TrendReq(hl="en-US", tz=0, retries=0, backoff_factor=0,
                  requests_args={"headers": {"User-Agent": "media-impacts-app/0.2"}})

    def _one_call(geo_val, timeframe_val):
        qs = all_q[:5]
        py.build_payload(qs, timeframe=timeframe_val, geo=geo_val)
        df = py.interest_over_time()
        if df is None or df.empty:
            return None, None
        if "isPartial" in df.columns:
            df = df.drop(columns=["isPartial"])
        fcols = [c for c in df.columns if c in film_qs]
        ocols = [c for c in df.columns if c in outcome_qs]
        if not fcols or not ocols:
            return None, None
        film_s = df[fcols].mean(axis=1).resample("W-MON").mean().rename("film_interest")
        out_s  = df[ocols].mean(axis=1).resample("W-MON").mean().rename("outcome")
        if film_s.empty or out_s.empty:
            return None, None
        return film_s, out_s

    if len(all_q) <= 5:
        fs, os_ = _one_call(geo, timeframe)
        if fs is not None: return fs, os_
        if timeframe != "today 5-y":
            fs, os_ = _one_call(geo, "today 5-y")
            if fs is not None: return fs, os_
        if geo:
            fs, os_ = _one_call("", "today 5-y")
            if fs is not None: return fs, os_
    fs = cached_trends(film_qs, geo, timeframe).rename("film_interest")
    os_ = cached_trends(outcome_qs, geo, timeframe).rename("outcome")
    return fs, os_

# =====================================================
# -------------- Paper 'results' access ---------------
# =====================================================
GITHUB_OWNER = "hsflabstanford"
GITHUB_REPO  = "media-impacts"
RESULTS_PATH = "results"

@st.cache_data(ttl=3600, show_spinner=False)
def github_list_results() -> Dict[str, str]:
    """Return {name: raw_url} for CSVs under results/ from the official repo (best-effort)."""
    api = f"https://api.github.com/repos/{GITHUB_OWNER}/{GITHUB_REPO}/contents/{RESULTS_PATH}"
    headers = {"User-Agent": "media-impacts-app/0.2"}
    out = {}
    try:
        r = requests.get(api, headers=headers, timeout=20)
        r.raise_for_status()
        for item in r.json():
            if item.get("type") == "file" and item.get("name", "").lower().endswith(".csv"):
                out[item["name"]] = item.get("download_url") or item.get("url")
        # Also walk one level of subfolders (some repos store per-run folders)
        for item in r.json():
            if item.get("type") == "dir":
                r2 = requests.get(item["url"], headers=headers, timeout=20)
                if r2.ok:
                    for it in r2.json():
                        if it.get("type") == "file" and it.get("name", "").lower().endswith(".csv"):
                            out[it["name"]] = it.get("download_url") or it.get("url")
    except Exception:
        pass
    return out

def infer_role_from_name(name: str) -> str:
    n = name.lower()
    if "all" in n and "doc" in n and "table" in n: return "all_docs_table"
    if "ci" in n and "assoc" in n: return "ci_assoc_contemp_lagged"
    if "ci" in n and "bin" in n and "ksu" in n: return "ci_ksu_lags"
    if "ci" in n and "bin" in n: return "ci_its"
    if "benchmark" in n: return "benchmarking_timeseries"
    if "media" in n and "plot" in n: return "media_plots"
    return "other"

def try_read_csv(url_or_bytes) -> Optional[pd.DataFrame]:
    try:
        if isinstance(url_or_bytes, (bytes, bytearray, io.BytesIO)):
            return pd.read_csv(io.BytesIO(url_or_bytes if isinstance(url_or_bytes, (bytes, bytearray)) else url_or_bytes.getvalue()))
        return pd.read_csv(url_or_bytes)
    except Exception:
        return None

def render_forest_plot(df: pd.DataFrame, label_col: str, effect_col: str, lo_col: str, hi_col: str, title: str) -> plt.Figure:
    Z = df[[label_col, effect_col, lo_col, hi_col]].dropna().copy()
    Z = Z.sort_values(effect_col)
    y = np.arange(len(Z))
    fig = plt.figure(figsize=(10, max(3, 0.28*len(Z)+1.5)))
    plt.hlines(y, Z[lo_col], Z[hi_col], linewidth=2)
    plt.plot(Z[effect_col], y, "o")
    plt.axvline(0, linestyle="--", alpha=0.6)
    plt.yticks(y, Z[label_col])
    plt.xlabel("Effect size (with 95% CI)")
    plt.title(title); plt.tight_layout()
    return fig

# =====================================================
# ----------------------- TABS ------------------------
# =====================================================
tab1, tab2 = st.tabs(["🆕 New Title (Quick ITS)", "📚 Study set (Paper)"])

# ---------------------- TAB 1 ------------------------
with tab1:
    with st.container():
        st.markdown('<div class="center-card">', unsafe_allow_html=True)
        with st.form("center_form"):
            # Minimal inputs; clearer label with example
            film_input = st.text_input(
                "Project title (e.g., Eating Our Way to Extinction)",
                value="",
                placeholder="Type a film/series title…",
            )
            intervention = st.date_input(
                "Date of release",
                value=dt.date(2021, 9, 30),
                help="Approximate is OK. We align to Mondays internally."
            )

            with st.expander("Advanced options", expanded=False):
                st.caption("What should we measure?")
                outcome_input = st.text_input(
                    "Outcome topic (default shown)",
                    value="Plant-based diet",
                    placeholder="e.g., vegan diet, meat consumption, food waste",
                    help="This is the outcome/idea/behavior you want to track."
                )

                st.caption("Where should we measure attention?")
                data_source = st.selectbox(
                    "Measurement source",
                    ["Wikipedia pageviews (recommended)", "Google Trends (index 0–100)"],
                    index=0,
                    help="Wikipedia is stable; Google Trends can rate-limit but shows relative interest."
                )
                auto_fallback = st.checkbox(
                    "If Google Trends fails, automatically use Wikipedia",
                    value=True
                )

                if "Google" in data_source:
                    colA, colB = st.columns(2)
                    with colA:
                        geo = st.selectbox("Region", ["", "US", "GB", "TH"], index=1, help="Blank = worldwide")
                    with colB:
                        timeframe = st.selectbox("Timeframe", ["today 5-y", "today 12-m", "today 3-m"], index=0)
                    show_regions = st.checkbox("Show interest by region (Trends only)", value=False)
                else:
                    geo = ""; timeframe = "today 5-y"; show_regions = False

                st.caption("Wikipedia settings")
                colC, colD, colE = st.columns(3)
                with colC:
                    start_daily = st.text_input("Start (YYYYMMDD)", "20190101")
                with colD:
                    end_daily = st.text_input("End (YYYYMMDD)", dt.date.today().strftime("%Y%m%d"))
                with colE:
                    lang = st.text_input("Language", "en")

                st.caption("Extras")
                multi_dates = st.text_input(
                    "Compare several dates (comma-separated, optional)",
                    "2021-09-30, 2022-07-01"
                )
                batch_file = st.file_uploader("Batch mode CSV (optional). Columns: film,outcome,intervention,source,geo",
                                            type=["csv"])

            run = st.form_submit_button("▶️ Run analysis")

        st.markdown("</div>", unsafe_allow_html=True)

    # Quick helpers
    st.session_state.setdefault("force_wiki", False)
    colx, coly = st.columns([1, 1])
    with colx:
        if st.button("🔄 Retry now"):
            st.rerun()
    with coly:
        if st.button("🧭 Switch to Wikipedia and retry"):
            st.session_state["force_wiki"] = True
            st.rerun()

    # Results area
    right = st.container()

    if run:
        try:
            with right, st.spinner("Fetching data and running ITS…"):
                use_wiki = ("Wikipedia" in data_source) or st.session_state.get("force_wiki", False)

                if not use_wiki:
                    film_qs = [q.strip() for q in (film_input or "").split(",")]
                    outcome_qs = [q.strip() for q in (outcome_input or "").split(",")]
                    try:
                        film_s, out_s = trends_pair(film_qs, outcome_qs, geo, timeframe)
                    except Exception as e:
                        if auto_fallback:
                            st.warning(f"Google Trends failed ({e}). Falling back to Wikipedia pageviews.")
                            film_s = cached_wiki(film_input, start_daily, end_daily, lang=lang).rename("film_interest")
                            out_s  = cached_wiki(outcome_input, start_daily, end_daily, lang=lang).rename("outcome")
                        else:
                            raise
                else:
                    film_s = cached_wiki(film_input, start_daily, end_daily, lang=lang).rename("film_interest")
                    out_s  = cached_wiki(outcome_input, start_daily, end_daily, lang=lang).rename("outcome")

                # Model
                model, df, figs, metrics = quick_its(out_s, film_s, intervention.isoformat())

                # ---------------- Results (clear wording) ----------------
                st.subheader("Results")

                def _fmt_int(x):
                    try:
                        if x is None or (isinstance(x, float) and np.isnan(x)): return "—"
                        return f"{int(round(x)):,}"
                    except Exception:
                        return str(x)

                def _fmt_float(x, d=3):
                    try:
                        if x is None or (isinstance(x, float) and np.isnan(x)): return "—"
                        return f"{float(x):.{d}f}"
                    except Exception:
                        return str(x)

                sig_level = 0.05
                p_post     = float(model.pvalues.get("post", 1.0))
                p_timepost = float(model.pvalues.get("time_post", 1.0))

                pre_avg = metrics.get("pre_mean", np.nan)
                pct = np.nan
                if pre_avg not in (None, 0) and not (isinstance(pre_avg, float) and np.isnan(pre_avg)):
                    try:
                        pct = 100.0 * metrics["level_change"] / pre_avg
                    except Exception:
                        pct = np.nan

                level_line = (
                    f"- **Immediate lift at release:** {_fmt_int(metrics['level_change'])} "
                    f"(95% CI {_fmt_int(metrics['level_ci_lo'])} to {_fmt_int(metrics['level_ci_hi'])}). "
                    f"{'This change is statistically significant.' if p_post < sig_level else 'This change is not statistically significant.'}"
                )

                if np.isnan(pct):
                    context_line = "- **Context:** We couldn’t compute a % of the pre-release average (too little or no pre-release data)."
                else:
                    context_line = f"- **Context:** about **{_fmt_float(pct, 1)}%** of a typical pre-release week."

                slope_val = metrics["slope_change_per_week"]
                slope_dir = "increased" if slope_val > 0 else "decreased" if slope_val < 0 else "did not change"
                slope_line = (
                    f"- **After release, the weekly trend {slope_dir} by** {_fmt_float(slope_val, 3)} per week "
                    f"(95% CI {_fmt_float(metrics['slope_ci_lo'], 3)} to {_fmt_float(metrics['slope_ci_hi'], 3)}). "
                    f"{'Significant.' if p_timepost < sig_level else 'Not significant.'}"
                )

                st.markdown(
                    f"""
**At the release week ({intervention}):**  
{level_line}  
{context_line}

**After the release:**  
{slope_line}
""")
                st.caption("Notes: “Significant” means p<0.05. Positive numbers = more attention; negative = less.")

                # Core charts
                st.pyplot(figs[0]); st.pyplot(figs[1])

                # Narrative explainer
                from scipy.optimize import curve_fit
                with st.expander("What these charts mean (plain language)", expanded=True):
                    design = sm.add_constant(df[["time","post","time_post","film_lag1"]])
                    pred    = model.predict(design)
                    cf = design.copy(); cf["post"]=0; cf["time_post"]=0
                    pred_cf = model.predict(cf)
                    after = df.index >= metrics["t0_monday"]
                    excess = (df["outcome"] - pred_cf).where(after, 0.0)
                    cum_excess = float(excess.cumsum().iloc[-1])
                    weeks_positive = int((excess > 0).sum())
                    peak_val = float(df["outcome"].max()); peak_when = df["outcome"].idxmax().date()

                    # Half-life (best-effort)
                    half_life_text = "n/a"
                    try:
                        def _exp_decay(t, A, k, c): return A*np.exp(-k*t) + c
                        pos = after & (excess > 0)
                        if pos.sum() >= 6:
                            y = excess[pos].values; t = np.arange(len(y))
                            (A, k, c), _ = curve_fit(_exp_decay, t, y, p0=[max(y),0.1,0.0], maxfev=20000)
                            hl = (np.log(2)/k) if k>0 else np.nan
                            if np.isfinite(hl): half_life_text = f"{hl:.1f} weeks"
                    except Exception:
                        pass

                    regions_text = ""
                    if ("Wikipedia" not in data_source) and not st.session_state.get("force_wiki", False):
                        try:
                            reg = cached_regions((film_input or "").split(",")[0].strip(), geo=geo, timeframe=timeframe)
                            qcol = (film_input or "").split(",")[0].strip()
                            top3 = reg.sort_values(qcol, ascending=False).head(3)["region"].tolist()
                            if top3: regions_text = f" Top regions: {', '.join(top3)}."
                        except Exception:
                            pass

                    sig = lambda p: "statistically significant" if p < 0.05 else "not statistically significant"
                    readable = (
                        f"**Headline** — Around **{intervention}**, attention to **{outcome_input}** changed.\n\n"
                        f"**Immediate jump:** {_fmt_int(metrics['level_change'])} "
                        f"({('%.1f' % (pct))+'%' if not np.isnan(pct) else 'n/a'} of pre), {sig(p_post)}.\n"
                        f"**Trend change:** {_fmt_float(metrics['slope_change_per_week'], 3)}/week, {sig(p_timepost)}.\n"
                        f"**Cumulative lift:** ~{_fmt_int(cum_excess)} units over {weeks_positive} weeks.\n"
                        f"**Peak:** {_fmt_int(peak_val)} on {peak_when}. **Half-life:** {half_life_text}.\n"
                        f"**Source:** {'Google Trends' if (('Wikipedia' not in data_source) and not st.session_state.get('force_wiki', False)) else 'Wikipedia pageviews'}."
                        f"{regions_text}\n"
                        f"*Correlations only; other events may also drive interest.*"
                    )
                    st.markdown(readable)
                    st.download_button("⬇️ Download summary (.txt)", readable.encode("utf-8"), file_name="readable_summary.txt")

                # Coefficients
                coef = (model.params.to_frame("coef")
                        .join(model.bse.to_frame("stderr"))
                        .join(model.pvalues.to_frame("pval")))
                st.subheader("Model coefficients")
                st.dataframe(coef.style.format({"coef":"{:.4f}","stderr":"{:.4f}","pval":"{:.4f}"}))

                # Downloads
                merged = df.copy(); merged.index.name = "week_start"
                st.download_button("⬇️ Coefficients (CSV)", coef.to_csv().encode(), file_name="its_coefficients.csv")
                st.download_button("⬇️ Weekly series (CSV)", merged.to_csv().encode(), file_name="weekly_series.csv")

                # ZIP bundle
                buf = io.BytesIO()
                with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
                    for i, fig in enumerate(figs, start=1):
                        fbuf = io.BytesIO(); fig.savefig(fbuf, format="png", dpi=200, bbox_inches="tight")
                        z.writestr(f"figure_{i}.png", fbuf.getvalue())
                    z.writestr("its_coefficients.csv", coef.to_csv().encode())
                    z.writestr("weekly_series.csv", merged.to_csv().encode())
                st.download_button("⬇️ Report bundle (ZIP)", buf.getvalue(), file_name="its_report_bundle.zip")

                # PowerPoint
                # (optional; comment out if you don't want pptx dependency)
                try:
                    from pptx import Presentation
                    from pptx.util import Inches
                    def make_ppt(figs, film, outcome, intervention, metrics) -> bytes:
                        prs = Presentation()
                        slide = prs.slides.add_slide(prs.slide_layouts[0])
                        slide.shapes.title.text = "Media Impacts — Quick ITS"
                        slide.placeholders[1].text = f"Title: {film}\nOutcome: {outcome}\nDate: {intervention}"
                        slide = prs.slides.add_slide(prs.slide_layouts[1]); slide.shapes.title.text = "Results"
                        body = slide.shapes.placeholders[1].text_frame
                        body.text = (f"Immediate change: {metrics['level_change']:.0f} "
                                     f"({metrics['level_change_pct_of_pre']:.1f}% of pre)\n"
                                     f"Slope change/wk: {metrics['slope_change_per_week']:.3f}\n"
                                     f"95% CIs — level: {metrics['level_ci_lo']:.0f}…{metrics['level_ci_hi']:.0f}; "
                                     f"slope: {metrics['slope_ci_lo']:.3f}…{metrics['slope_ci_hi']:.3f}")
                        for fig in figs:
                            slide2 = prs.slides.add_slide(prs.slide_layouts[5])
                            img = io.BytesIO(); fig.savefig(img, format="png", dpi=200, bbox_inches="tight"); img.seek(0)
                            slide2.shapes.add_picture(img, Inches(0.5), Inches(1.0), width=Inches(8.5))
                        b = io.BytesIO(); prs.save(b); b.seek(0); return b.getvalue()
                    st.download_button("⬇️ PowerPoint (.pptx)",
                                       make_ppt(figs, film_input, outcome_input, intervention, metrics),
                                       file_name="its_report.pptx")
                except Exception:
                    pass

                # -------- Performance visuals with captions --------
                st.subheader("Performance visuals")
                design = sm.add_constant(df[["time","post","time_post","film_lag1"]])
                pred    = model.predict(design)
                cf = design.copy(); cf["post"]=0; cf["time_post"]=0
                pred_cf = model.predict(cf)
                after_mask = df.index >= metrics["t0_monday"]
                excess = (df["outcome"] - pred_cf).where(after_mask, 0.0)

                # 1) Actual vs Counterfactual
                fig = plt.figure(figsize=(10,4))
                plt.plot(df.index, df["outcome"], label="Actual outcome")
                plt.plot(df.index, pred_cf, label="Counterfactual (no intervention)")
                plt.fill_between(df.index[after_mask], pred_cf[after_mask], df["outcome"][after_mask], alpha=0.25)
                plt.axvline(metrics["t0_monday"], linestyle="--", label="Intervention")
                plt.title("Actual vs Counterfactual"); plt.legend(); plt.tight_layout()
                st.pyplot(fig)
                st.caption("Blue = actual; orange = ‘no-release’ path. Shaded area = extra attention after the date.")

                # 2) Weekly excess
                fig = plt.figure(figsize=(10,3))
                plt.bar(df.index[after_mask], excess[after_mask])
                plt.axvline(metrics["t0_monday"], linestyle="--")
                plt.title("Weekly Excess vs Counterfactual"); plt.tight_layout()
                st.pyplot(fig)
                st.caption("Weekly lift above (or below) baseline after the intervention.")

                # 3) Cumulative excess
                cum_excess = excess.cumsum()
                fig = plt.figure(figsize=(10,3))
                plt.plot(df.index, cum_excess)
                plt.axvline(metrics["t0_monday"], linestyle="--")
                plt.title("Cumulative Excess Interest"); plt.tight_layout()
                st.pyplot(fig)
                st.caption("Running total of extra attention vs counterfactual.")

                # 4) Event-time average (±26 weeks)
                weeks_from = ((df.index - metrics["t0_monday"]).days // 7).astype(int)
                window = weeks_from.between(-26, 26)
                es = df.loc[window, ["outcome"]].copy()
                es["k"] = weeks_from[window]
                es_avg = es.groupby("k")["outcome"].mean()
                fig = plt.figure(figsize=(10,3))
                plt.plot(es_avg.index, es_avg.values)
                plt.axvline(0, linestyle="--")
                plt.title("Event-time average (weeks relative to intervention)")
                plt.xlabel("Weeks from intervention"); plt.tight_layout()
                st.pyplot(fig)
                st.caption("Centered view of lead-up and after-effect around the release week (0).")

                # 5) Decay fit (best-effort)
                from scipy.optimize import curve_fit
                def _exp_decay(t, A, k, c): return A * np.exp(-k * t) + c
                pos = after_mask & (excess > 0)
                if pos.sum() >= 6:
                    y = excess[pos].values; t = np.arange(len(y))
                    try:
                        (A, k, c), _ = curve_fit(_exp_decay, t, y, p0=[max(y), 0.1, 0.0], maxfev=20000)
                        hl = (np.log(2) / k) if k > 0 else np.nan
                        fig = plt.figure(figsize=(10,3))
                        plt.plot(df.index[pos], y, label="Observed excess")
                        plt.plot(df.index[pos], _exp_decay(t, A, k, c), label=f"Decay fit (half-life ≈ {hl:.1f} wks)")
                        plt.axvline(metrics["t0_monday"], linestyle="--")
                        plt.title("Excess decay after intervention"); plt.legend(); plt.tight_layout()
                        st.pyplot(fig)
                        st.caption("How fast the lift fades. Half-life ≈ weeks to drop by half.")
                    except Exception as _e:
                        st.info(f"Decay fit skipped: {_e}")
                        st.caption("Couldn’t fit a stable decay curve this time.")

                # 6) Placebo check
                def _fit_at(date_monday):
                    d = df.copy()
                    d["post"] = (d.index >= date_monday).astype(int)
                    d["time_post"] = d["time"] * d["post"]
                    Xp = sm.add_constant(d[["time","post","time_post","film_lag1"]])
                    return sm.OLS(d["outcome"], Xp).fit(cov_type="HAC", cov_kwds={"maxlags":4})
                pre_idx = df.index[df.index < metrics["t0_monday"]][8:]
                n_placebo = int(min(50, max(0, len(pre_idx)-8))); placebo_effects = []
                if n_placebo >= 10:
                    rng = np.random.default_rng(42)
                    picks = np.sort(rng.choice(pre_idx, size=n_placebo, replace=False))
                    for d0 in picks:
                        res = _fit_at(d0)
                        if "post" in res.params:
                            placebo_effects.append(float(res.params["post"]))
                    if placebo_effects:
                        fig = plt.figure(figsize=(10,3))
                        plt.hist(placebo_effects, bins=20, alpha=0.8)
                        plt.axvline(metrics["level_change"], color="red", linestyle="--", label="Observed level change")
                        plt.title("Placebo distribution (pre-period dates)"); plt.legend(); plt.tight_layout()
                        st.pyplot(fig)
                        st.caption("Grey bars = fake pre-dates. Red line = your jump; further right ⇒ less likely by chance.")

                # 7) Regions (Trends only)
                if ("Wikipedia" not in data_source) and not st.session_state.get("force_wiki", False) and show_regions:
                    try:
                        st.subheader("Interest by region (Top 15)")
                        reg = cached_regions((film_input or "").split(",")[0].strip(), geo=geo, timeframe=timeframe)
                        qcol = (film_input or "").split(",")[0].strip()
                        top = reg.sort_values(qcol, ascending=False).head(15)
                        fig = plt.figure(figsize=(10,5))
                        plt.barh(top["region"][::-1], top[qcol][::-1])
                        plt.title("Where interest is highest"); plt.tight_layout()
                        st.pyplot(fig)
                        st.caption("Locations with the strongest search interest for the project (Google Trends).")
                    except Exception as e:
                        st.info(f"Region breakdown unavailable: {e}")

                # Compare dates (optional)
                if multi_dates.strip():
                    dates = [d.strip() for d in multi_dates.split(",") if d.strip()]
                    try:
                        rows = []
                        for d in dates:
                            m2, df2, _, met2 = quick_its(out_s, film_s, d)
                            rows.append({
                                "intervention": d,
                                "level_change": met2["level_change"],
                                "level_%_pre": met2["level_change_pct_of_pre"],
                                "level_p": float(m2.pvalues.get("post", np.nan)),
                                "slope_change": met2["slope_change_per_week"],
                                "slope_p": float(m2.pvalues.get("time_post", np.nan)),
                            })
                        table = pd.DataFrame(rows).sort_values("level_%_pre", ascending=False)
                        st.subheader("Date comparison")
                        st.dataframe(table.style.format({
                            "level_change": "{:.0f}", "level_%_pre": "{:.1f}%", "level_p": "{:.3f}",
                            "slope_change": "{:.2f}", "slope_p": "{:.3f}"
                        }))
                        st.download_button("⬇️ Download comparison (CSV)", table.to_csv(index=False).encode(),
                                           file_name="its_date_comparison.csv")
                    except Exception as e:
                        st.warning(f"Could not run date comparison: {e}")

        except Exception as e:
            st.error(f"Error: {e}")
            st.stop()

# ---------------------- TAB 2 ------------------------
with tab2:
    st.write("Use this tab to view prepared results from the paper (tables, forest plots, benchmarking).")
    src = st.radio("Data source", ["Auto-fetch from official GitHub (best effort)", "Upload CSVs"], index=0)

    # Placeholders for discovered dataframes
    found: Dict[str, pd.DataFrame] = {}
    ci_candidates = {}

    if src.startswith("Auto"):
        with st.spinner("Looking for CSVs in the repo’s results/ folder…"):
            listing = github_list_results()
        if not listing:
            st.warning("Couldn’t read the repo’s results listing. You can switch to **Upload CSVs** below.")
        else:
            # Try to pull key tables if present
            roles = {}
            for name, url in listing.items():
                role = infer_role_from_name(name)
                roles.setdefault(role, []).append((name, url))

            def fetch_first(role: str) -> Optional[pd.DataFrame]:
                for name, url in roles.get(role, []):
                    df = try_read_csv(url)
                    if df is not None and not df.empty:
                        return df
                return None

            found["all_docs_table"] = fetch_first("all_docs_table")
            found["benchmarking"]   = fetch_first("benchmarking_timeseries")
            # forest plot candidates (store entire mapping and let user pick columns)
            for key in ("ci_assoc_contemp_lagged", "ci_its", "ci_ksu_lags"):
                if roles.get(key):
                    # load first, but also keep the url for re-load if user switches
                    name, url = roles[key][0]
                    df = try_read_csv(url)
                    if df is not None and not df.empty:
                        ci_candidates[key] = (name, df)

    else:
        st.info("Upload one or more CSV files produced by the repo’s notebooks (e.g., create-all-docs-table, create-ci-... ).")
        uploads = st.file_uploader("Drop CSVs here", type=["csv"], accept_multiple_files=True)
        if uploads:
            for up in uploads:
                name = up.name
                role = infer_role_from_name(name)
                df   = try_read_csv(up.getvalue())
                if df is not None and not df.empty:
                    if role.startswith("ci_"):
                        ci_candidates[role] = (name, df)
                    elif role == "benchmarking_timeseries":
                        found["benchmarking"] = df
                    elif role == "all_docs_table":
                        found["all_docs_table"] = df

    # ---- All-Docs table
    st.subheader("All-docs table")
    if found.get("all_docs_table") is not None:
        df_all = found["all_docs_table"].copy()
        # light clean
        for col in df_all.columns:
            if df_all[col].dtype == object:
                try:
                    df_all[col] = df_all[col].str.strip()
                except Exception:
                    pass
        c1, c2 = st.columns([2,1])
        with c1:
            st.dataframe(df_all, use_container_width=True)
        with c2:
            st.download_button("⬇️ Download (CSV)", df_all.to_csv(index=False).encode(), file_name="all_docs_table.csv")
    else:
        st.info("No “all-docs” table found yet. If you have it locally, upload the CSV.")

    # ---- Forest plots
    st.subheader("Forest plots")
    if ci_candidates:
        choices = {"Association/Contemporaneous/Lagged": "ci_assoc_contemp_lagged",
                   "Interrupted Time Series (ITS)": "ci_its",
                   "KSU Lagged": "ci_ksu_lags"}
        label = st.selectbox("Choose result set", list(choices.keys()))
        key   = choices[label]
        if key in ci_candidates:
            name, df_ci = ci_candidates[key]
            st.caption(f"Loaded: **{name}**")
            # Guess columns
            cols = df_ci.columns.str.lower()
            # Common patterns
            label_guess = next((c for c in df_ci.columns if "title" in c.lower() or "film" in c.lower() or "name" in c.lower()), df_ci.columns[0])
            eff_guess   = next((c for c in df_ci.columns if "effect" in c.lower() or "estimate" in c.lower() or "coef" in c.lower()), None)
            lo_guess    = next((c for c in df_ci.columns if "ci_lo" in c.lower() or "lower" in c.lower() or "lo" in c.lower()), None)
            hi_guess    = next((c for c in df_ci.columns if "ci_hi" in c.lower() or "upper" in c.lower() or "hi" in c.lower()), None)

            col1, col2, col3, col4 = st.columns(4)
            with col1:
                label_col = st.selectbox("Label column", df_ci.columns.tolist(), index=df_ci.columns.get_loc(label_guess))
            with col2:
                effect_col = st.selectbox("Effect column", df_ci.columns.tolist(), index=(df_ci.columns.get_loc(eff_guess) if eff_guess else 0))
            with col3:
                lo_col = st.selectbox("Lower CI", df_ci.columns.tolist(), index=(df_ci.columns.get_loc(lo_guess) if lo_guess else 0))
            with col4:
                hi_col = st.selectbox("Upper CI", df_ci.columns.tolist(), index=(df_ci.columns.get_loc(hi_guess) if hi_guess else 0))

            fig = render_forest_plot(df_ci, label_col, effect_col, lo_col, hi_col, f"{label} — effect sizes")
            st.pyplot(fig)
            st.download_button("⬇️ Download data (CSV)", df_ci.to_csv(index=False).encode(), file_name=f"{key}.csv")
        else:
            st.info("Couldn’t load the selected forest data.")
    else:
        st.info("No forest CSVs detected yet. Try auto-fetch again or upload the files from the visualization notebooks.")

    # ---- Benchmarking time series
    st.subheader("Benchmarking time series")
    if found.get("benchmarking") is not None:
        df_b = found["benchmarking"].copy()
        st.dataframe(df_b.head(50), use_container_width=True)
        # Try best-effort line plot (needs 'date' + several series)
        date_col = next((c for c in df_b.columns if "date" in c.lower()), None)
        if date_col:
            try:
                df_b[date_col] = pd.to_datetime(df_b[date_col])
                fig = plt.figure(figsize=(10,4))
                for c in df_b.columns:
                    if c == date_col: continue
                    try:
                        plt.plot(df_b[date_col], pd.to_numeric(df_b[c], errors="coerce"), label=c)
                    except Exception:
                        pass
                plt.legend(); plt.title("Benchmarking (paper reproduction)"); plt.tight_layout()
                st.pyplot(fig)
            except Exception as e:
                st.info(f"Couldn’t plot benchmarking series: {e}")
        st.download_button("⬇️ Download benchmarking (CSV)", df_b.to_csv(index=False).encode(),
                           file_name="benchmarking_timeseries.csv")
    else:
        st.info("No benchmarking CSV detected. Upload the CSV produced by the notebook if you want this view.")
